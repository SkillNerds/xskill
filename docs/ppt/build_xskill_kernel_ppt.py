"""生成《XSkill 可插拔技能生产管线重构设计》中文培训 PPT。

运行：
    python3.11 docs/ppt/build_xskill_kernel_ppt.py
"""

from __future__ import annotations

from pathlib import Path
from textwrap import dedent

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "XSkill_可插拔技能生产管线重构设计_算法团队培训版.pptx"

W = 13.333
H = 7.5

FONT = "Noto Sans CJK SC"
FONT_MEDIUM = "Noto Sans CJK SC Medium"
FONT_MONO = "Noto Sans Mono CJK SC"

NAVY = "0B1020"
NAVY_2 = "151B2F"
NAVY_3 = "202A46"
WHITE = "F7F9FC"
TEXT = "111827"
TEXT_2 = "394150"
MUTED = "AAB4C8"
MUTED_2 = "6B7280"
LIGHT = "F3F6FA"
LIGHT_2 = "E9EEF6"
BORDER = "D6DEEA"
TEAL = "27D3B2"
TEAL_DARK = "0F8F7B"
VIOLET = "7B61FF"
BLUE = "4DA3FF"
AMBER = "F4B740"
RED = "EF6A74"
GREEN = "20B978"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = TEXT,
    font: str = FONT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    line_spacing: float = 1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    runs: list[tuple[str, dict]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.05
    for text, opts in runs:
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", FONT)
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = rgb(opts.get("color", color))
    return box


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 1,
    radius: bool = False,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = BORDER,
    width: float = 1.5,
    dash: bool = False,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = 2
    return line


def add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = MUTED_2,
    width: float = 1.6,
):
    add_line(slide, x1, y1, x2 - 0.10, y2, color=color, width=width)
    tri = add_shape(
        slide,
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        x2 - 0.13,
        y2 - 0.07,
        0.15,
        0.14,
        fill=color,
    )
    tri.rotation = 90
    return tri


def add_title(
    slide,
    section: str,
    title: str,
    subtitle: str = "",
    *,
    dark: bool = False,
):
    color = WHITE if dark else TEXT
    sub_color = MUTED if dark else MUTED_2
    add_text(slide, section.upper(), 0.55, 0.27, 2.1, 0.30, size=10, color=TEAL, bold=True)
    add_text(slide, title, 0.55, 0.62, 12.1, 0.60, size=27, color=color, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.22, 11.9, 0.34, size=12.5, color=sub_color)


def add_footer(slide, page: int, *, dark: bool = False, source: str = ""):
    line_color = NAVY_3 if dark else BORDER
    text_color = MUTED if dark else MUTED_2
    add_line(slide, 0.55, 7.13, 12.78, 7.13, color=line_color, width=0.7)
    if source:
        add_text(slide, source, 0.58, 7.19, 10.9, 0.18, size=7.7, color=text_color)
    add_text(
        slide,
        f"{page:02d}",
        12.05,
        7.17,
        0.65,
        0.20,
        size=8,
        color=text_color,
        align=PP_ALIGN.RIGHT,
    )


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    accent: str = TEAL,
    dark: bool = False,
    title_size: float = 16,
    body_size: float = 11.5,
):
    fill = NAVY_2 if dark else WHITE
    border = NAVY_3 if dark else BORDER
    title_color = WHITE if dark else TEXT
    body_color = MUTED if dark else TEXT_2
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=fill, line=border, radius=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.07, h, fill=accent, radius=True)
    add_text(slide, title, x + 0.25, y + 0.18, w - 0.42, 0.35, size=title_size, color=title_color, bold=True)
    add_text(slide, body, x + 0.25, y + 0.63, w - 0.43, h - 0.78, size=body_size, color=body_color, line_spacing=1.12)


def add_chip(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: str = LIGHT_2,
    color: str = TEXT_2,
    border: str | None = None,
    size: float = 10,
):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.32, fill=fill, line=border, radius=True)
    add_text(
        slide,
        text,
        x + 0.07,
        y + 0.01,
        w - 0.14,
        0.28,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_code_block(
    slide,
    code: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str = "",
    size: float = 10.5,
):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=NAVY, line=NAVY_3, radius=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.30, fill=NAVY_2, radius=True)
    for i, c in enumerate((RED, AMBER, GREEN)):
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.18 + i * 0.20, y + 0.11, 0.08, 0.08, fill=c)
    if label:
        add_text(slide, label, x + 0.85, y + 0.04, w - 1.02, 0.20, size=8.5, color=MUTED)
    box = add_text(
        slide,
        dedent(code).strip(),
        x + 0.18,
        y + 0.41,
        w - 0.36,
        h - 0.52,
        size=size,
        color="DCE6F5",
        font=FONT_MONO,
        line_spacing=0.96,
    )
    box.text_frame.margin_left = Inches(0.02)
    return box


def add_number_badge(slide, n: int, x: float, y: float, *, color: str = TEAL):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, 0.34, 0.34, fill=color)
    add_text(
        slide,
        str(n),
        x,
        y + 0.01,
        0.34,
        0.28,
        size=10,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.OVAL, 9.9, -1.0, 4.5, 4.5, fill=NAVY_2)
    add_shape(slide, MSO_SHAPE.OVAL, 10.9, -0.05, 2.7, 2.7, fill=VIOLET)
    add_shape(slide, MSO_SHAPE.OVAL, 11.55, 0.58, 1.4, 1.4, fill=TEAL)
    add_chip(slide, "ARCHITECTURE / TRAINING", 0.68, 0.65, 2.20, fill=NAVY_3, color=TEAL)
    add_text(slide, "XSkill 可插拔技能生产管线", 0.72, 1.68, 10.8, 0.75, size=36, color=WHITE, bold=True)
    add_text(slide, "重构设计", 0.72, 2.48, 6.0, 0.70, size=36, color=TEAL, bold=True)
    add_text(
        slide,
        "面向算法团队的接入、热配置与评测指南",
        0.76,
        3.48,
        8.2,
        0.42,
        size=18,
        color=MUTED,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.75, 4.33, 4.35, 0.04, fill=TEAL)
    add_text(
        slide,
        "核心命题：平台负责轨迹与 Skill 生命周期；算法内核只负责从标准输入生成标准产物。",
        0.76,
        4.67,
        9.5,
        0.76,
        size=17,
        color=WHITE,
        bold=True,
        line_spacing=1.12,
    )
    add_text(slide, "2026.07 · 设计提案 / 培训版", 0.76, 6.56, 4.3, 0.28, size=10.5, color=MUTED)
    add_text(slide, "参考：OpenMines BaseDispatcher 扩展方式", 8.0, 6.56, 4.55, 0.28, size=10.5, color=MUTED, align=PP_ALIGN.RIGHT)
    return slide


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "01 / 结论", "把平台做稳定，把算法做可替换", "这次重构的目标不是“多一个配置项”，而是建立可长期演进的算法边界。")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 1.82, 12.05, 1.28, fill=NAVY, radius=True)
    add_rich_text(
        slide,
        [
            ("XSkill = ", {"color": MUTED, "size": 19}),
            ("稳定外壳", {"color": WHITE, "size": 23, "bold": True}),
            ("  +  ", {"color": MUTED, "size": 19}),
            ("统一内核接口", {"color": TEAL, "size": 23, "bold": True}),
            ("  +  ", {"color": MUTED, "size": 19}),
            ("版本化运行时", {"color": WHITE, "size": 23, "bold": True}),
            ("  +  ", {"color": MUTED, "size": 19}),
            ("平台统一评价", {"color": AMBER, "size": 23, "bold": True}),
        ],
        0.98,
        2.14,
        11.4,
        0.55,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    cards = [
        ("接口稳定", "算法只接收标准请求，只返回 SkillArtifact；不碰上传、用户统计和下发数据库。", TEAL),
        ("包级注册", "生产环境用 Python package entry point；开发环境保留 module:Class 快速调试。", VIOLET),
        ("无损切换", "前端提交配置修订；新任务原子切换，在途任务继续固定旧版本，可一键回滚。", BLUE),
        ("评价归平台", "统一数据集、统一指标、统一裁决；算法可上报诊断，但不能决定自己是否胜出。", AMBER),
    ]
    for i, (title, body, color) in enumerate(cards):
        add_card(slide, 0.65 + i * 3.05, 3.52, 2.78, 2.35, title, body, accent=color, title_size=16, body_size=11.2)
        add_number_badge(slide, i + 1, 2.82 + i * 3.05, 3.68, color=color)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 6.23, 12.05, 0.52, fill="E8FBF7", line="BDEDE3", radius=True)
    add_text(
        slide,
        "推荐起点：先把现有 AtomTask 管线封装成 builtin-atomtask 内核，行为不变；再开放第三方内核。",
        0.89,
        6.35,
        11.55,
        0.25,
        size=11.5,
        color=TEAL_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 2)
    return slide


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "02 / 边界", "谁负责什么：边界必须一眼能看懂", "算法团队只对“如何生成 Skill”负责；平台对数据、用户和发布负责。")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.62, 1.78, 5.73, 4.72, fill="F7F9FC", line=BORDER, radius=True)
    add_chip(slide, "XSkill 平台外壳（稳定）", 0.96, 2.02, 2.42, fill=NAVY, color=WHITE)
    platform_items = [
        ("轨迹接入", "上传、去重、脱敏、用户归属"),
        ("运行编排", "队列、重试、配置修订、资源配额"),
        ("Skill 生命周期", "仓库、版本、灰度、下发、回滚"),
        ("统计与评价", "成本、时延、质量、线上反馈、报表"),
    ]
    for i, (t, b) in enumerate(platform_items):
        y = 2.58 + i * 0.88
        add_shape(slide, MSO_SHAPE.OVAL, 0.98, y + 0.05, 0.26, 0.26, fill=TEAL)
        add_text(slide, t, 1.42, y, 1.45, 0.30, size=13.5, color=TEXT, bold=True)
        add_text(slide, b, 2.74, y + 0.01, 3.14, 0.40, size=11.3, color=TEXT_2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.96, 1.78, 5.74, 4.72, fill=NAVY, line=NAVY_3, radius=True)
    add_chip(slide, "算法内核（可替换）", 7.29, 2.02, 2.06, fill=TEAL, color=NAVY)
    kernel_items = [
        ("拆分", "trajectory → AtomTask / 其他中间表示"),
        ("聚类与路由", "合并、去重、候选 Skill 决策"),
        ("生成", "SKILL.md、scripts、references"),
        ("算法诊断", "阶段事件、内部指标、lineage"),
    ]
    for i, (t, b) in enumerate(kernel_items):
        y = 2.58 + i * 0.88
        add_shape(slide, MSO_SHAPE.OVAL, 7.31, y + 0.05, 0.26, 0.26, fill=VIOLET if i < 3 else AMBER)
        add_text(slide, t, 7.75, y, 1.45, 0.30, size=13.5, color=WHITE, bold=True)
        add_text(slide, b, 9.08, y + 0.01, 3.14, 0.40, size=11.3, color=MUTED)

    add_shape(slide, MSO_SHAPE.CHEVRON, 5.94, 3.17, 1.43, 1.28, fill=WHITE, line=BORDER)
    add_text(slide, "标准\n契约", 6.20, 3.44, 0.72, 0.52, size=13.5, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 2.61, 6.70, 8.12, 0.33, fill="FFF5DD", radius=True)
    add_text(
        slide,
        "禁止跨界：内核不能直接更新用户统计、checkout Skill 分支或调用下发接口。",
        2.78,
        6.73,
        7.78,
        0.23,
        size=10.5,
        color="9A6411",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 3)
    return slide


def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "03 / 现状审计", "现有能力可复用，主要问题是“编排与实现绑死”", "以下判断来自当前仓库 runner、TaskAgent、TaskClusterAgent、SkillEditAgent 与 canary 实现。")

    x_positions = [0.67, 3.05, 5.43, 7.81, 10.19]
    labels = [
        ("DirectoryWatcher", "扫描 / 调度", BLUE),
        ("TaskAgent", "轨迹拆分", TEAL),
        ("TaskClusterAgent", "归类 / buffer", VIOLET),
        ("SkillEditAgent", "生成 / 提交", AMBER),
        ("Canary", "评价 / 晋级", GREEN),
    ]
    for i, (name, sub, color) in enumerate(labels):
        add_shape(slide, MSO_SHAPE.RECTANGLE, x_positions[i], 2.14, 1.89, 1.00, fill=WHITE, line=BORDER, radius=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x_positions[i], 2.14, 1.89, 0.08, fill=color, radius=True)
        add_text(slide, name, x_positions[i] + 0.12, 2.42, 1.65, 0.25, size=11.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x_positions[i] + 0.12, 2.76, 1.65, 0.22, size=9.7, color=MUTED_2, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_arrow(slide, x_positions[i] + 1.91, 2.63, x_positions[i + 1] - 0.04, 2.63, color=MUTED_2)

    add_text(slide, "当前耦合点", 0.70, 3.63, 2.0, 0.35, size=16, color=RED, bold=True)
    issues = [
        ("直接构造类", "Watcher 直接 import / 调用具体 Agent"),
        ("共享内部状态", "config、store、git、DB、目录结构同时暴露"),
        ("没有版本固定", "任务无法回答“由哪个内核 + 哪版配置生成”"),
        ("替换粒度过细", "改其中一段容易破坏跨阶段隐含约定"),
    ]
    for i, (t, b) in enumerate(issues):
        x = 0.68 + (i % 2) * 3.06
        y = 4.08 + (i // 2) * 0.93
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 2.78, 0.69, fill="FFF0F2", line="F5C7CD", radius=True)
        add_text(slide, t, x + 0.16, y + 0.10, 0.96, 0.23, size=11, color="A43642", bold=True)
        add_text(slide, b, x + 1.05, y + 0.10, 1.55, 0.44, size=9.8, color="7D3A42")

    add_text(slide, "应保留的资产", 7.10, 3.63, 2.0, 0.35, size=16, color=GREEN, bold=True)
    assets = [
        ("AtomTask 数据模型", "输入证据可追溯"),
        ("候选 buffer 与阈值", "增量积累机制成熟"),
        ("main / staging", "已经具备版本与灰度基础"),
        ("UX / cost / model 统计", "可直接成为统一评测数据"),
    ]
    for i, (t, b) in enumerate(assets):
        x = 7.08 + (i % 2) * 2.77
        y = 4.08 + (i // 2) * 0.93
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 2.50, 0.69, fill="EAF9F3", line="BFE9D7", radius=True)
        add_text(slide, t, x + 0.15, y + 0.10, 1.20, 0.23, size=10.5, color="137453", bold=True)
        add_text(slide, b, x + 1.20, y + 0.10, 1.13, 0.40, size=9.6, color="2E6B58")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 6.42, 12.00, 0.44, fill=NAVY, radius=True)
    add_text(slide, "重构策略：把现有整条算法链先封装为一个默认内核，不立即重写算法逻辑。", 0.90, 6.52, 11.55, 0.22, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(
        slide,
        4,
        source="代码依据：src/xskill/pipeline/runner.py · agents/task_agent.py · task_cluster_agent.py · skill_edit_agent.py · canary.py",
    )
    return slide


def slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "04 / 参考模式", "从 OpenMines 学什么：统一基类 + 配置选择实现", "OpenMines 用 BaseDispatcher 固定方法集合，用 JSON 选择算法；XSkill 需要在此基础上补齐生产治理。")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.62, 1.78, 5.55, 4.91, fill=NAVY, line=NAVY_3, radius=True)
    add_chip(slide, "OpenMines", 0.92, 2.04, 1.22, fill=TEAL, color=NAVY)
    add_code_block(
        slide,
        """
        class BaseDispatcher:
            def give_init_order(self, truck, mine): ...
            def give_haul_order(self, truck, mine): ...
            def give_back_order(self, truck, mine): ...

        class NaiveDispatcher(BaseDispatcher):
            def give_init_order(self, truck, mine):
                return 0

        # config.json
        {"dispatcher": {"type": ["NaiveDispatcher"]}}
        """,
        0.92,
        2.55,
        4.93,
        3.48,
        label="BaseDispatcher + config",
        size=9.8,
    )
    add_text(slide, "值得保留", 0.94, 6.18, 1.1, 0.23, size=10.5, color=TEAL, bold=True)
    add_text(slide, "• 基类定义最小接口  • 算法与框架分工清楚  • 配置决定运行实现", 2.02, 6.17, 3.82, 0.30, size=9.8, color=MUTED)

    add_text(slide, "翻译到 XSkill", 6.62, 1.94, 2.3, 0.34, size=18, color=TEXT, bold=True)
    comparisons = [
        ("BaseDispatcher", "SkillGenerationKernel", "固定输入 / 输出，不固定内部步骤"),
        ("算法目录", "独立 wheel 包", "允许独立版本、依赖与发布节奏"),
        ("JSON type", "entry point + 配置修订", "显式包名、版本与 schema"),
        ("运行后对比产量", "平台离线 / 影子 / 灰度评价", "统一数据与裁决规则"),
    ]
    for i, (left, right, note) in enumerate(comparisons):
        y = 2.52 + i * 0.91
        add_shape(slide, MSO_SHAPE.RECTANGLE, 6.62, y, 1.54, 0.54, fill=LIGHT_2, radius=True)
        add_text(slide, left, 6.70, y + 0.13, 1.37, 0.22, size=9.6, color=TEXT_2, bold=True, align=PP_ALIGN.CENTER)
        add_arrow(slide, 8.23, y + 0.27, 8.69, y + 0.27, color=VIOLET)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 8.78, y, 1.86, 0.54, fill="EEEAFE", line="D9D0FF", radius=True)
        add_text(slide, right, 8.87, y + 0.13, 1.68, 0.22, size=9.6, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, note, 10.83, y + 0.05, 1.72, 0.43, size=9.2, color=TEXT_2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.62, 6.26, 5.93, 0.43, fill="FFF5DD", line="F4DDAA", radius=True)
    add_text(
        slide,
        "不照搬：生产环境不扫描目录猜类名；注册必须显式、可校验、可审计。",
        6.80,
        6.35,
        5.58,
        0.23,
        size=10.4,
        color="8A5A10",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(
        slide,
        5,
        source="参考：https://github.com/370025263/openmines（README / openmines/src/dispatcher.py / cli/run.py）",
    )
    return slide


def slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(slide, "05 / 总体架构", "目标架构：控制、执行、数据与评价分层", "对外保持一个 XSkill；对内把算法内核作为版本化插件运行。", dark=True)

    layers = [
        (1.85, "控制层", "前端配置  ·  API  ·  插件目录  ·  配置修订  ·  权限 / 审计", BLUE),
        (3.01, "执行层", "任务路由  ·  Runtime Manager  ·  builtin / team-a / team-b kernels", VIOLET),
        (5.10, "数据与评价", "Trajectory Store  ·  Artifact Store  ·  Skill Repo  ·  Evaluator  ·  Distributor", TEAL),
    ]
    for y, label, body, color in layers:
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.64, y, 12.04, 0.79 if y != 3.01 else 1.55, fill=NAVY_2, line=NAVY_3, radius=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.64, y, 0.12, 0.79 if y != 3.01 else 1.55, fill=color, radius=True)
        add_text(slide, label, 0.97, y + 0.22, 1.55, 0.30, size=15, color=color, bold=True)
        if y != 3.01:
            add_text(slide, body, 2.48, y + 0.23, 9.74, 0.27, size=12, color=WHITE)
        else:
            add_text(slide, "任务路由", 2.45, y + 0.21, 1.10, 0.30, size=12, color=MUTED, bold=True)
            add_arrow(slide, 3.51, y + 0.36, 4.08, y + 0.36, color=MUTED)
            add_shape(slide, MSO_SHAPE.RECTANGLE, 4.14, y + 0.12, 2.12, 0.48, fill=NAVY_3, line=VIOLET, radius=True)
            add_text(slide, "Runtime Manager", 4.28, y + 0.24, 1.84, 0.22, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            add_arrow(slide, 6.29, y + 0.36, 6.82, y + 0.36, color=MUTED)
            kernels = [
                ("builtin-atomtask", TEAL),
                ("team-a@1.3", VIOLET),
                ("team-b@0.9", AMBER),
            ]
            for i, (name, c) in enumerate(kernels):
                x = 6.92 + i * 1.77
                add_shape(slide, MSO_SHAPE.RECTANGLE, x, y + 0.12, 1.52, 0.48, fill=NAVY_3, line=c, radius=True)
                add_text(slide, name, x + 0.07, y + 0.24, 1.38, 0.22, size=9.2, color=c, bold=True, align=PP_ALIGN.CENTER)
            add_text(
                slide,
                "每个 run 固定：kernel_name + kernel_version + config_revision",
                2.45,
                y + 0.94,
                9.55,
                0.27,
                size=11.5,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )

    add_arrow(slide, 6.67, 2.67, 6.67, 2.96, color=BLUE)
    add_arrow(slide, 6.67, 4.63, 6.67, 5.05, color=TEAL)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.66, 6.34, 9.99, 0.44, fill=NAVY_3, line=VIOLET, radius=True)
    add_text(
        slide,
        "关键原则：平台通过 SDK 服务向内核提供能力；内核不反向依赖 xskill 内部模块。",
        1.92,
        6.44,
        9.48,
        0.22,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 6, dark=True)
    return slide


def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "06 / 数据契约", "替换的是算法，不能替换输入输出语义", "统一契约让不同管线可以比较、重放、回滚，也让前端知道如何展示配置与结果。")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 2.02, 3.23, 3.76, fill="EAF3FF", line="C8DCF5", radius=True)
    add_chip(slide, "KernelRunRequest", 0.98, 2.27, 1.72, fill=BLUE, color=NAVY)
    request_items = [
        "run_id / tenant / user",
        "trajectory_refs（只读引用）",
        "kernel + version",
        "config_revision",
        "resource_budget / deadline",
        "trace_context",
    ]
    for i, text in enumerate(request_items):
        add_shape(slide, MSO_SHAPE.OVAL, 1.03, 2.93 + i * 0.42, 0.10, 0.10, fill=BLUE)
        add_text(slide, text, 1.25, 2.87 + i * 0.42, 2.18, 0.25, size=10.7, color=TEXT_2)

    add_shape(slide, MSO_SHAPE.HEXAGON, 4.55, 2.75, 4.17, 2.36, fill=NAVY, line=VIOLET)
    add_text(slide, "SkillGenerationKernel", 5.04, 3.29, 3.19, 0.35, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "内部步骤完全自由", 5.29, 3.78, 2.69, 0.28, size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "TaskAgent 不是公共接口\nAtomTask 可以是某个内核的内部表示", 5.12, 4.16, 3.05, 0.58, size=10.3, color=MUTED, align=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 9.40, 2.02, 3.23, 3.76, fill="ECFAF6", line="C7EADF", radius=True)
    add_chip(slide, "KernelRunResult", 9.73, 2.27, 1.67, fill=TEAL, color=NAVY)
    result_items = [
        "SkillArtifact[]",
        "SKILL.md + 辅助文件",
        "lineage（来源轨迹 / atom）",
        "metrics（诊断，不作裁决）",
        "events / warnings",
        "可复现的运行元数据",
    ]
    for i, text in enumerate(result_items):
        add_shape(slide, MSO_SHAPE.OVAL, 9.78, 2.93 + i * 0.42, 0.10, 0.10, fill=TEAL)
        add_text(slide, text, 10.00, 2.87 + i * 0.42, 2.18, 0.25, size=10.7, color=TEXT_2)

    add_arrow(slide, 3.91, 3.92, 4.47, 3.92, color=BLUE, width=2)
    add_arrow(slide, 8.75, 3.92, 9.33, 3.92, color=TEAL, width=2)

    add_card(
        slide,
        0.66,
        6.08,
        3.76,
        0.72,
        "只读输入",
        "内核不能修改原轨迹与用户元数据。",
        accent=BLUE,
        title_size=11.5,
        body_size=9.4,
    )
    add_card(
        slide,
        4.78,
        6.08,
        3.76,
        0.72,
        "不可变产物",
        "先写 Artifact Store，通过门禁后再发布。",
        accent=VIOLET,
        title_size=11.5,
        body_size=9.4,
    )
    add_card(
        slide,
        8.89,
        6.08,
        3.76,
        0.72,
        "全链路追溯",
        "每个 Skill 能回答来源与生成版本。",
        accent=TEAL,
        title_size=11.5,
        body_size=9.4,
    )
    add_footer(slide, 7)
    return slide


def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "07 / SDK", "统一抽象接口：算法团队只实现四个入口", "接口要小、数据结构要强；不要把现有 TaskAgent / ClusterAgent 类直接变成公共 API。")

    add_code_block(
        slide,
        """
        class SkillGenerationKernel(ABC):
            @classmethod
            def manifest(cls) -> KernelManifest: ...

            @classmethod
            def config_schema(cls) -> Mapping[str, Any]: ...

            def validate_config(self) -> None: ...

            def run(
                self,
                request: KernelRunRequest,
                services: KernelServices,
            ) -> KernelRunResult: ...
        """,
        0.66,
        1.87,
        7.10,
        4.88,
        label="xskill-kernel-sdk / contract.py",
        size=12.3,
    )

    add_text(slide, "为什么只暴露 run？", 8.15, 1.95, 3.8, 0.36, size=18, color=TEXT, bold=True)
    reasons = [
        ("允许算法自由组织阶段", "有的内核是三段 Agent，有的可能是端到端模型或工作流引擎。", TEAL),
        ("减少跨版本锁定", "平台不依赖 AtomTask、candidate buffer 等某一种内部实现。", VIOLET),
        ("评价以产物为中心", "统一比较最终 Skill、lineage、成本和时延，而不是比较内部日志格式。", AMBER),
    ]
    for i, (t, b, c) in enumerate(reasons):
        add_card(slide, 8.14, 2.49 + i * 1.24, 4.45, 1.02, t, b, accent=c, title_size=12.3, body_size=9.7)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.14, 6.29, 4.45, 0.46, fill=NAVY, radius=True)
    add_text(slide, "KernelServices 只给受控能力：读轨迹、调用模型、写临时产物、发事件。", 8.34, 6.39, 4.03, 0.25, size=9.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)
    return slide


def slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "08 / 包规范", "推荐交付方式：独立 Python 包 + entry point", "脚本路径适合本地调试；生产环境必须可版本化、可安装、可校验。")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 1.83, 4.27, 4.95, fill=LIGHT, line=BORDER, radius=True)
    add_chip(slide, "包目录", 0.96, 2.08, 0.93, fill=NAVY, color=WHITE)
    tree = """
    xskill-kernel-team-a/
    ├── pyproject.toml
    ├── README.md
    ├── team_a_kernel/
    │   ├── kernel.py       ← SDK 适配层
    │   ├── config.py       ← 配置模型 / schema
    │   ├── algo_core/      ← 算法包
    │   └── resources/
    └── tests/
        ├── test_contract.py
        └── fixtures/
    """
    add_text(slide, dedent(tree).strip(), 1.02, 2.67, 3.64, 3.35, size=11.8, color=TEXT_2, font=FONT_MONO, line_spacing=1.05)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.00, 6.21, 3.56, 0.36, fill="E8FBF7", radius=True)
    add_text(slide, "kernel.py 薄；核心算法仍由算法团队原包维护。", 1.15, 6.28, 3.26, 0.20, size=9.7, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)

    add_code_block(
        slide,
        """
        [project]
        name = "xskill-kernel-team-a"
        version = "1.3.0"
        dependencies = [
          "xskill-kernel-sdk>=1,<2",
          "team-a-algo==4.2.1",
        ]

        [project.entry-points."xskill.skillgen_kernels"]
        team-a = "team_a_kernel.kernel:TeamAKernel"
        """,
        5.32,
        1.84,
        7.35,
        3.18,
        label="pyproject.toml",
        size=11.2,
    )

    add_text(slide, "两种加载模式", 5.35, 5.35, 2.0, 0.30, size=16, color=TEXT, bold=True)
    add_card(
        slide,
        5.34,
        5.77,
        3.45,
        0.92,
        "开发模式",
        "module:Class；本地快速修改和调试。",
        accent=BLUE,
        title_size=12,
        body_size=9.6,
    )
    add_card(
        slide,
        9.13,
        5.77,
        3.54,
        0.92,
        "生产模式",
        "package + version + entry point；安装后进入注册表。",
        accent=TEAL,
        title_size=12,
        body_size=9.6,
    )
    add_footer(slide, 9)
    return slide


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(slide, "09 / 热切换", "“热更换”定义：服务不停，新任务切换，在途任务不变", "不要在同一 Python 进程里 reload 已运行的模块；依赖、缓存和全局状态都可能失控。", dark=True)

    steps = [
        ("保存草稿", "前端编辑配置", BLUE),
        ("校验", "schema / 兼容性", TEAL),
        ("预热", "安装 + smoke", VIOLET),
        ("激活", "原子更新路由", AMBER),
        ("观察 / 回滚", "影子或灰度", GREEN),
    ]
    x0 = 0.73
    for i, (title, sub, color) in enumerate(steps):
        x = x0 + i * 2.48
        add_shape(slide, MSO_SHAPE.OVAL, x, 2.21, 0.48, 0.48, fill=color)
        add_text(slide, str(i + 1), x, 2.29, 0.48, 0.24, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x - 0.12, 2.88, 1.58, 0.28, size=13.5, color=WHITE, bold=True)
        add_text(slide, sub, x - 0.12, 3.25, 1.72, 0.26, size=10, color=MUTED)
        if i < len(steps) - 1:
            add_arrow(slide, x + 0.56, 2.45, x + 2.28, 2.45, color=NAVY_3, width=2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.79, 4.03, 5.78, 1.43, fill=NAVY_2, line=BLUE, radius=True)
    add_chip(slide, "旧任务", 1.09, 4.28, 0.85, fill=BLUE, color=NAVY)
    add_text(slide, "run-1007", 2.19, 4.28, 1.15, 0.25, size=12, color=WHITE, bold=True)
    add_text(slide, "team-a@1.2.0 · config rev 41", 2.18, 4.70, 3.75, 0.28, size=12.5, color=MUTED)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.09, 5.10, 5.06, 0.08, fill=BLUE, radius=True)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.83, 4.03, 5.78, 1.43, fill=NAVY_2, line=TEAL, radius=True)
    add_chip(slide, "新任务", 7.13, 4.28, 0.85, fill=TEAL, color=NAVY)
    add_text(slide, "run-1008", 8.23, 4.28, 1.15, 0.25, size=12, color=WHITE, bold=True)
    add_text(slide, "team-a@1.3.0 · config rev 42", 8.22, 4.70, 3.75, 0.28, size=12.5, color=MUTED)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 7.13, 5.10, 5.06, 0.08, fill=TEAL, radius=True)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.64, 6.08, 10.05, 0.55, fill="2B2134", line=RED, radius=True)
    add_text(
        slide,
        "切换失败时只回滚 active revision 指针；不覆盖旧修订，不修改在途任务记录。",
        1.89,
        6.22,
        9.52,
        0.25,
        size=11.2,
        color="FFC2C8",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 10, dark=True)
    return slide


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "10 / 前端与配置", "前端不是直接改 YAML，而是提交不可变的配置修订", "配置修订 = 内核身份 + 参数 + 发布策略；后端完成校验、预热和激活。")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 1.84, 6.17, 4.96, fill=WHITE, line=BORDER, radius=True)
    add_text(slide, "技能生产内核", 0.94, 2.10, 2.2, 0.34, size=17, color=TEXT, bold=True)
    fields = [
        ("内核", "team-a", 2.65),
        ("版本", "1.3.0", 3.35),
        ("配置修订", "rev 42（草稿）", 4.05),
        ("发布方式", "Shadow → Canary 20%", 4.75),
    ]
    for label, value, y in fields:
        add_text(slide, label, 0.96, y, 1.15, 0.28, size=10.5, color=MUTED_2, bold=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 2.06, y - 0.06, 4.35, 0.43, fill=LIGHT, line=BORDER, radius=True)
        add_text(slide, value, 2.24, y + 0.02, 3.99, 0.23, size=10.8, color=TEXT, bold=True)
    add_text(slide, "内核参数（由 config_schema 自动生成表单）", 0.96, 5.48, 4.20, 0.28, size=10.5, color=MUTED_2, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.96, 5.86, 5.45, 0.60, fill=NAVY, radius=True)
    add_text(slide, '{"min_cluster_size": 3, "batch_size": 8}', 1.14, 6.05, 5.08, 0.24, size=10.7, color=TEAL, font=FONT_MONO)
    add_chip(slide, "校验并保存", 5.19, 6.53, 1.23, fill=TEAL, color=NAVY)

    add_text(slide, "后端事务", 7.23, 1.98, 2.0, 0.34, size=18, color=TEXT, bold=True)
    api_steps = [
        ("POST /kernel-configs", "生成 rev 42，状态 DRAFT"),
        ("POST /kernel-configs/42/validate", "schema + SDK + capability"),
        ("POST /kernel-configs/42/prewarm", "安装包、构造实例、smoke"),
        ("POST /kernel-configs/42/activate", "CAS 更新 active_revision"),
        ("POST /kernel-configs/41/activate", "回滚仍是一次激活"),
    ]
    for i, (api, result) in enumerate(api_steps):
        y = 2.53 + i * 0.77
        add_number_badge(slide, i + 1, 7.24, y, color=TEAL if i < 4 else RED)
        add_text(slide, api, 7.78, y - 0.01, 4.43, 0.25, size=10.2, color=TEXT, bold=True, font=FONT_MONO)
        add_text(slide, result, 7.78, y + 0.29, 4.45, 0.23, size=9.5, color=TEXT_2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 7.23, 6.45, 5.34, 0.36, fill="FFF5DD", radius=True)
    add_text(slide, "密钥只保存引用（secret_ref），不能进入算法配置正文。", 7.40, 6.52, 5.02, 0.22, size=9.8, color="8A5A10", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)
    return slide


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "11 / 评价", "评价必须与内核解耦：同题、同预算、同裁决", "内核可以报告内部诊断；最终质量判断由 XSkill Evaluator 对标准产物执行。")

    stages = [
        ("合同门禁", "schema / 文件安全 / lineage 完整", BLUE),
        ("离线基准", "固定轨迹集 / 多次复跑 / 成本预算", VIOLET),
        ("影子运行", "同一输入双跑，不向用户下发", AMBER),
        ("线上灰度", "真实反馈、最小样本、自动晋级 / 回滚", GREEN),
    ]
    for i, (title, body, color) in enumerate(stages):
        x = 0.68 + i * 3.05
        w = 2.74 - i * 0.12
        add_shape(slide, MSO_SHAPE.CHEVRON, x, 2.02, w, 1.03, fill=color)
        add_text(slide, title, x + 0.16, 2.24, w - 0.41, 0.25, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.20, 2.58, w - 0.48, 0.31, size=9.2, color=NAVY, align=PP_ALIGN.CENTER)

    add_text(slide, "统一评分卡", 0.70, 3.54, 2.0, 0.35, size=17, color=TEXT, bold=True)
    metrics = [
        ("产物有效率", "硬门禁", "Skill schema / 安全 / 可安装", BLUE),
        ("任务覆盖率", "质量", "多少可复用意图被正确提炼", TEAL),
        ("触发准确率", "质量", "Precision / Recall / F1", VIOLET),
        ("去重与稳定性", "质量", "近义 Skill 重复率 / 多次复跑一致性", AMBER),
        ("用户体验", "线上", "真实 UX、使用后撤销 / 重做", GREEN),
        ("成本与时延", "效率", "每轨迹 Token、P95、失败率", RED),
    ]
    for i, (metric, tag, desc, color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = 0.69 + col * 4.13
        y = 4.02 + row * 1.03
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 3.82, 0.83, fill=LIGHT, line=BORDER, radius=True)
        add_chip(slide, tag, x + 0.16, y + 0.15, 0.70, fill=color, color=NAVY, size=8.5)
        add_text(slide, metric, x + 1.01, y + 0.12, 1.44, 0.25, size=11.5, color=TEXT, bold=True)
        add_text(slide, desc, x + 1.01, y + 0.42, 2.56, 0.22, size=9.3, color=TEXT_2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.69, 6.28, 12.00, 0.48, fill=NAVY, radius=True)
    add_text(
        slide,
        "裁决顺序：先过硬门禁 → 再比较质量 → 质量相近时选成本更低者；不能只看一个加权总分。",
        0.95,
        6.40,
        11.47,
        0.24,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 12)
    return slide


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "12 / 接入流程", "算法团队交付一个新内核，只需要 5 步", "下面从“已有算法包”出发，演示如何接入 XSkill，而不是要求算法团队重写业务逻辑。")

    steps = [
        ("01", "保留算法包", "splitter / clusterer / writer 继续独立维护", BLUE),
        ("02", "编写适配类", "继承 SkillGenerationKernel，实现 4 个入口", TEAL),
        ("03", "声明配置", "提供 config_schema、默认值与不兼容说明", VIOLET),
        ("04", "注册与测试", "entry point + 合同测试 + 固定样例", AMBER),
        ("05", "提交评测", "离线 → Shadow → Canary，通过后激活", GREEN),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        x = 0.68 + i * 2.50
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 2.12, 2.15, 3.24, fill=WHITE, line=BORDER, radius=True)
        add_text(slide, num, x + 0.18, 2.34, 0.72, 0.43, size=26, color=color, bold=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.18, 2.96, 1.78, 0.05, fill=color)
        add_text(slide, title, x + 0.18, 3.30, 1.80, 0.33, size=15, color=TEXT, bold=True)
        add_text(slide, body, x + 0.18, 3.88, 1.76, 0.92, size=10.8, color=TEXT_2, line_spacing=1.16)
        if i < len(steps) - 1:
            add_arrow(slide, x + 2.17, 3.69, x + 2.42, 3.69, color=MUTED_2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.68, 5.79, 5.80, 0.76, fill="EAF3FF", line="C9DDF5", radius=True)
    add_text(slide, "算法团队负责", 0.92, 5.97, 1.26, 0.27, size=12, color=BLUE, bold=True)
    add_text(slide, "算法、配置 schema、合同测试、算法说明", 2.10, 5.98, 3.94, 0.25, size=10.8, color=TEXT_2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.82, 5.79, 5.86, 0.76, fill="ECFAF6", line="C4EADF", radius=True)
    add_text(slide, "平台团队负责", 7.07, 5.97, 1.26, 0.27, size=12, color=TEAL_DARK, bold=True)
    add_text(slide, "安装、路由、资源、评价、发布、回滚、下发", 8.25, 5.98, 3.95, 0.25, size=10.8, color=TEXT_2)
    add_footer(slide, 13)
    return slide


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "13 / Demo", "Demo 项目：核心算法与 XSkill 适配层分开", "PPT 同目录附带可运行代码：docs/ppt/xskill_kernel_demo。")

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.66, 1.84, 5.08, 4.90, fill=NAVY, line=NAVY_3, radius=True)
    add_chip(slide, "xskill_kernel_demo/", 0.97, 2.10, 1.98, fill=TEAL, color=NAVY)
    tree = """
    ├── xskill_kernel_sdk.py     # 稳定接口（演示版）
    ├── config.yaml              # 选择内核 + 修订
    ├── pyproject.toml           # entry point
    ├── run_demo.py              # 本地宿主
    └── demo_skillgen/
        ├── algo_core.py         # 算法团队已有包
        └── kernel.py            # 薄适配层
    """
    add_text(slide, dedent(tree).strip(), 1.02, 2.75, 4.32, 2.85, size=11.7, color="DCE6F5", font=FONT_MONO, line_spacing=1.09)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.97, 5.97, 4.44, 0.45, fill=NAVY_3, line=VIOLET, radius=True)
    add_text(slide, "运行：python3.11 run_demo.py", 1.16, 6.08, 4.06, 0.23, size=10.5, color=TEAL, font=FONT_MONO, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, "一次调用发生什么", 6.14, 1.99, 2.42, 0.34, size=18, color=TEXT, bold=True)
    flow = [
        ("标准请求", "KernelRunRequest", BLUE),
        ("适配层", "DemoAtomTaskKernel", VIOLET),
        ("算法包", "Splitter → Clusterer → Writer", AMBER),
        ("标准结果", "KernelRunResult", TEAL),
    ]
    for i, (title, sub, color) in enumerate(flow):
        y = 2.55 + i * 0.96
        add_shape(slide, MSO_SHAPE.RECTANGLE, 6.14, y, 5.93, 0.66, fill=LIGHT, line=BORDER, radius=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 6.14, y, 0.09, 0.66, fill=color, radius=True)
        add_text(slide, title, 6.47, y + 0.17, 1.38, 0.25, size=12, color=color, bold=True)
        add_text(slide, sub, 7.91, y + 0.17, 3.77, 0.25, size=11.2, color=TEXT, font=FONT_MONO if i != 2 else FONT)
        if i < len(flow) - 1:
            add_arrow(slide, 9.10, y + 0.69, 9.10, y + 0.91, color=MUTED_2)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.14, 6.46, 5.93, 0.29, fill="E8FBF7", radius=True)
    add_text(slide, "算法包完全不知道上传、用户统计、Skill 下发在哪里。", 6.34, 6.49, 5.54, 0.20, size=9.8, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 14)
    return slide


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_title(slide, "14 / Demo 代码", "算法内核怎么写：继承基类，再调用自己的算法包", "kernel.py 只做协议适配，不把算法逻辑搬进 XSkill。", dark=True)

    add_code_block(
        slide,
        """
        class DemoAtomTaskKernel(SkillGenerationKernel):
            def __init__(self, config):
                self.config = dict(config)
                self.splitter = SimpleSplitter()
                self.clusterer = KeywordClusterer()
                self.writer = MarkdownSkillWriter()

            @classmethod
            def manifest(cls):
                return KernelManifest(
                    name="demo-atomtask",
                    version="1.0.0",
                    capabilities=("batch", "lineage"),
                )

            @classmethod
            def config_schema(cls):
                return DEMO_CONFIG_SCHEMA
        """,
        0.66,
        1.83,
        6.01,
        4.97,
        label="demo_skillgen/kernel.py · 1/2",
        size=10.8,
    )
    add_code_block(
        slide,
        """
        def run(self, request, services):
            atoms = []
            for traj in request.trajectories:
                text = services.read_text(traj.path)
                atoms += self.splitter.split(
                    traj.trajectory_id, text
                )

            groups = self.clusterer.cluster(
                atoms,
                min_cluster_size=
                    self.config["min_cluster_size"],
            )
            artifacts = [
                SkillArtifact(
                    name=name,
                    content=self.writer.write(name, members),
                    source_trajectory_ids=source_ids(members),
                )
                for name, members in groups.items()
            ]
            return KernelRunResult(
                artifacts=artifacts,
                metrics={"atom_count": len(atoms)},
                lineage=build_lineage(artifacts),
            )
        """,
        6.94,
        1.83,
        5.73,
        4.97,
        label="demo_skillgen/kernel.py · 2/2",
        size=9.8,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.74, 6.35, 9.85, 0.42, fill=NAVY_3, line=TEAL, radius=True)
    add_text(
        slide,
        "注意：调用模型、读取轨迹、写临时产物都走 services；不要 import xskill.pipeline.runner。",
        2.02,
        6.45,
        9.28,
        0.22,
        size=10.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 15, dark=True)
    return slide


def slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "15 / Demo 配置", "配置决定选谁，包元数据决定怎么找到它", "开发时用 class_path；生产时用 package + version + entry point。")

    add_code_block(
        slide,
        """
        skill_generation:
          active:
            class_path: demo_skillgen.kernel:DemoAtomTaskKernel
            package: xskill-kernel-demo-atomtask
            version: 1.0.0
          config_revision: 42
          config:
            min_cluster_size: 1
          rollout:
            mode: shadow
            shadow_kernel: builtin-atomtask
        """,
        0.65,
        1.84,
        5.77,
        3.59,
        label="config.yaml",
        size=10.8,
    )
    add_code_block(
        slide,
        """
        [project]
        name = "xskill-kernel-demo-atomtask"
        version = "1.0.0"

        [project.entry-points."xskill.skillgen_kernels"]
        demo-atomtask =
          "demo_skillgen.kernel:DemoAtomTaskKernel"
        """,
        6.73,
        1.84,
        5.95,
        2.19,
        label="pyproject.toml",
        size=10.4,
    )
    add_code_block(
        slide,
        """
        $ python3.11 run_demo.py
        manifest: demo-atomtask@1.0.0
        events: atoms.ready count=2
        metrics: atom_count=2, skill_count=1
        skills: ['修复']
        lineage: {'修复': ['traj-demo']}
        """,
        6.73,
        4.30,
        5.95,
        2.29,
        label="smoke output",
        size=10.4,
    )

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.65, 5.77, 5.77, 0.82, fill=WHITE, line=BORDER, radius=True)
    add_text(slide, "切换到新算法", 0.91, 5.95, 1.37, 0.27, size=12.5, color=VIOLET, bold=True)
    add_text(slide, "只新增一个配置修订；旧配置仍保留，可重放、可回滚。", 2.24, 5.96, 3.79, 0.28, size=10.5, color=TEXT_2)
    add_footer(slide, 16)
    return slide


def slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "16 / 端到端", "从前端切换到 Skill 下发：一条可审计的链路", "每个关键动作都有状态、有版本、有回滚点。")

    lanes = [
        ("前端 / API", BLUE),
        ("Config Service", TEAL),
        ("Runtime Manager", VIOLET),
        ("Kernel Worker", AMBER),
        ("Evaluator / Repo", GREEN),
        ("Distributor", RED),
    ]
    x_positions = [0.52, 2.57, 4.62, 6.67, 8.72, 10.77]
    for i, (name, color) in enumerate(lanes):
        add_shape(slide, MSO_SHAPE.RECTANGLE, x_positions[i], 1.84, 1.73, 0.46, fill=color, radius=True)
        add_text(slide, name, x_positions[i] + 0.06, 1.96, 1.61, 0.22, size=9.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_line(slide, x_positions[i] + 0.865, 2.35, x_positions[i] + 0.865, 6.53, color=BORDER, width=0.8, dash=True)

    events = [
        (0, 1, 2.58, "保存 rev 42"),
        (1, 2, 3.13, "validate + prewarm"),
        (2, 0, 3.68, "返回可激活"),
        (0, 1, 4.23, "activate rev 42"),
        (1, 2, 4.78, "原子更新路由"),
        (2, 3, 5.33, "新 run 固定 1.3.0 / rev 42"),
        (3, 4, 5.88, "提交 Artifact + lineage"),
        (4, 5, 6.43, "门禁通过后下发"),
    ]
    colors = [BLUE, TEAL, VIOLET, AMBER, GREEN, RED]
    for source, target, y, label in events:
        x1 = x_positions[source] + 0.865
        x2 = x_positions[target] + 0.865
        direction = 1 if x2 > x1 else -1
        if direction == 1:
            add_arrow(slide, x1, y, x2, y, color=colors[source], width=1.7)
            tx = x1 + 0.10
        else:
            add_line(slide, x2 + 0.11, y, x1, y, color=colors[source], width=1.7)
            tri = add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x2 - 0.02, y - 0.07, 0.15, 0.14, fill=colors[source])
            tri.rotation = 270
            tx = x2 + 0.17
        add_shape(slide, MSO_SHAPE.RECTANGLE, tx, y - 0.20, abs(x2 - x1) - 0.22, 0.30, fill=WHITE)
        add_text(slide, label, tx + 0.04, y - 0.16, abs(x2 - x1) - 0.30, 0.20, size=8.6, color=TEXT_2, bold=True, align=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.20, 6.76, 10.92, 0.26, fill="E8FBF7", radius=True)
    add_text(slide, "状态真源：配置修订表 + run 记录 + Artifact 元数据；不要只依赖进程内对象。", 1.42, 6.79, 10.48, 0.18, size=9.3, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 17)
    return slide


def slide_18(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_title(slide, "17 / 落地路线", "先建立边界，再开放生态：四个阶段逐步迁移", "每个阶段都能独立验收，不要求一次性把当前管线重写完。")

    phases = [
        ("P0", "基线与合同", "冻结现有结果样例\n定义 Request / Result / Services\n建立合同测试", BLUE, "验收：现有 10 条轨迹可重放"),
        ("P1", "默认内核封装", "把现有三 Agent 链包进\nbuiltin-atomtask\nWatcher 只调 SDK", TEAL, "验收：行为与成本不回退"),
        ("P2", "插件与热配置", "entry point 注册\n配置修订 / 预热 / 激活\n前端配置表单", VIOLET, "验收：新任务切换、旧任务固定"),
        ("P3", "统一评价与开放", "离线基准 / Shadow / Canary\n包签名与资源限制\n算法团队自助发布", AMBER, "验收：自动晋级与回滚闭环"),
    ]
    for i, (p, title, body, color, gate) in enumerate(phases):
        x = 0.66 + i * 3.06
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.92, 2.77, 3.43, fill=WHITE, line=BORDER, radius=True)
        add_chip(slide, p, x + 0.18, 2.14, 0.55, fill=color, color=NAVY)
        add_text(slide, title, x + 0.18, 2.64, 2.34, 0.33, size=15, color=TEXT, bold=True)
        add_text(slide, body, x + 0.18, 3.25, 2.34, 1.14, size=10.8, color=TEXT_2, line_spacing=1.18)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.18, 4.68, 2.39, 0.42, fill=LIGHT, radius=True)
        add_text(slide, gate, x + 0.27, 4.78, 2.21, 0.23, size=9.1, color=color, bold=True, align=PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            add_arrow(slide, x + 2.79, 3.61, x + 3.00, 3.61, color=MUTED_2)

    add_text(slide, "今天需要拍板的 4 个决定", 0.68, 5.74, 3.1, 0.34, size=17, color=TEXT, bold=True)
    decisions = [
        ("1", "公共接口以整条管线 run 为粒度"),
        ("2", "生产插件只接受 wheel + entry point"),
        ("3", "热切换只影响新任务，所有 run 固定版本"),
        ("4", "评价与发布由平台统一控制"),
    ]
    for i, (n, text) in enumerate(decisions):
        x = 0.70 + i * 3.02
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 6.20, 2.72, 0.57, fill=NAVY, radius=True)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.14, 6.32, 0.27, 0.27, fill=TEAL)
        add_text(slide, n, x + 0.14, 6.35, 0.27, 0.18, size=8.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, text, x + 0.53, 6.30, 2.03, 0.31, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 18)
    return slide


def slide_19(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_chip(slide, "TAKEAWAY", 0.70, 0.68, 1.23, fill=NAVY_3, color=TEAL)
    add_text(slide, "算法团队只需记住三件事", 0.72, 1.45, 7.8, 0.64, size=32, color=WHITE, bold=True)

    takeaways = [
        ("01", "继承统一基类", "实现 manifest / config_schema / validate_config / run。", TEAL),
        ("02", "调用自己的算法包", "XSkill 适配层要薄，不依赖平台内部模块。", VIOLET),
        ("03", "交付可评价产物", "返回 SkillArtifact、lineage 和诊断；平台负责是否发布。", AMBER),
    ]
    for i, (num, title, body, color) in enumerate(takeaways):
        y = 2.55 + i * 1.15
        add_text(slide, num, 0.75, y, 0.64, 0.39, size=20, color=color, bold=True)
        add_text(slide, title, 1.52, y - 0.02, 2.55, 0.35, size=17, color=WHITE, bold=True)
        add_text(slide, body, 4.06, y + 0.01, 5.91, 0.30, size=12.2, color=MUTED)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 9.73, 1.37, 2.80, 4.78, fill=NAVY_2, line=NAVY_3, radius=True)
    add_text(slide, "接入检查", 10.03, 1.72, 1.86, 0.36, size=18, color=TEAL, bold=True)
    checklist = [
        "包有明确版本",
        "配置可校验",
        "不直接访问平台 DB",
        "结果有 lineage",
        "合同测试可重复",
        "支持失败与超时",
        "通过影子后再灰度",
    ]
    for i, item in enumerate(checklist):
        y = 2.34 + i * 0.48
        add_shape(slide, MSO_SHAPE.RECTANGLE, 10.04, y, 0.22, 0.22, fill=TEAL, radius=True)
        add_text(slide, "✓", 10.04, y - 0.01, 0.22, 0.18, size=8.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, item, 10.43, y - 0.01, 1.74, 0.24, size=10.3, color=WHITE)

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 6.40, 7.95, 0.49, fill=TEAL, radius=True)
    add_text(slide, "下一步：用 builtin-atomtask 完成 P0 / P1 原型，拿真实轨迹做等价性回归。", 0.98, 6.51, 7.42, 0.24, size=11.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Q & A", 10.17, 6.52, 2.17, 0.34, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 19, dark=True, source="配套 Demo：docs/ppt/xskill_kernel_demo")
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "XSkill 可插拔技能生产管线重构设计"
    prs.core_properties.subject = "算法团队接入、热配置与评测培训"
    prs.core_properties.author = "370025263"
    prs.core_properties.keywords = "XSkill, skill generation, plugin, pipeline, evaluation"
    prs.core_properties.comments = "基于 traj2skill 当前代码与 OpenMines BaseDispatcher 模式设计"

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    slide_15(prs)
    slide_16(prs)
    slide_17(prs)
    slide_18(prs)
    slide_19(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
